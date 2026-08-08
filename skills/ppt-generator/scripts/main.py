#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""PPT 生成工具 — 根据 JSON 数据生成 PowerPoint 演示文稿。

需要 python-pptx 库。如未安装会自动提示。
"""
import argparse
import json
import sys

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8")
    sys.stderr.reconfigure(encoding="utf-8")


def main():
    parser = argparse.ArgumentParser(description="PPT 生成工具")
    parser.add_argument("--slides", required=True, help="JSON 格式幻灯片数据")
    parser.add_argument("--output", required=True, help="输出 .pptx 文件路径")
    args = parser.parse_args()

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        print(json.dumps({"error": "需要 python-pptx 库，请运行: pip install python-pptx"}, ensure_ascii=False))
        sys.exit(1)

    try:
        data = json.loads(args.slides)
    except json.JSONDecodeError as e:
        print(json.dumps({"error": f"JSON 解析失败: {e}"}, ensure_ascii=False))
        sys.exit(1)

    slides_data = data.get("slides", [data] if isinstance(data, dict) else [])
    prs = Presentation()

    for slide_data in slides_data:
        title = slide_data.get("title", "")
        content = slide_data.get("content", "")
        bullets = slide_data.get("bullets", [])

        if content and not bullets:
            # 标题+内容页
            slide_layout = prs.slide_layouts[1]  # Title and Content
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title
            slide.placeholders[1].text = content
        elif bullets:
            # 要点页
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title
            tf = slide.placeholders[1].text_frame
            for i, bullet in enumerate(bullets):
                if i == 0:
                    tf.text = bullet
                else:
                    p = tf.add_paragraph()
                p.text = bullet
        else:
            # 纯标题页
            slide_layout = prs.slide_layouts[0]  # Title Slide
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title

    prs.save(args.output)
    print(json.dumps({"success": True, "slides_count": len(slides_data), "output_file": args.output}, ensure_ascii=False))


if __name__ == "__main__":
    main()
